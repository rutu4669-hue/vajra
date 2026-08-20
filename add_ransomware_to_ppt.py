#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_ransomware_content():
    # Load existing presentation
    prs = Presentation('DATE - 06-08-2026.pptx')
    
    # Add Ransomware Overview slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Ransomware Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "Ransomware is a type of malicious software that encrypts files and demands payment for decryption."
    
    p = content_frame.add_paragraph()
    p.text = "Key Characteristics:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    
    characteristics = [
        "• File encryption with strong algorithms",
        "• Ransom demand in cryptocurrency",
        "• Time-sensitive payment deadlines",
        "• Double extortion (encryption + data theft)",
        "• Targeted attacks on organizations",
        "• Ransomware-as-a-Service (RaaS) model"
    ]
    
    for char in characteristics:
        p = content_frame.add_paragraph()
        p.text = char
        p.font.size = Pt(16)
        p.level = 1
        p.space_before = Pt(6)
    
    # Add Ransomware Statistics slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Ransomware Statistics"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    stats = [
        "• Average ransom payment: $812,000 (2023)",
        "• 37% of organizations hit by ransomware",
        "• Average downtime: 21 days",
        "• 66% of victims lose data permanently",
        "• Healthcare is most targeted industry",
        "• 4,000+ ransomware attacks daily (2023)",
        "• Total ransomware damages: $20 billion (2023)",
        "• 90% of attacks originate from phishing"
    ]
    
    for stat in stats:
        p = content_frame.add_paragraph()
        p.text = stat
        p.font.size = Pt(18)
        p.space_before = Pt(8)
    
    # Add Major Ransomware Groups slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Major Ransomware Groups"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    groups = [
        "LockBit - Most active ransomware group (2023)",
        "BlackCat/ALPHV - Sophisticated RaaS platform",
        "Conti - Russian-speaking cybercrime group",
        "REvil/Sodinokibi - High-profile attacks",
        "DarkSide - Colonial Pipeline attack",
        "Hive - Healthcare sector targeting",
        "Karakurt - Data extortion specialist",
        "Lapsus$ - Teenage hacker group"
    ]
    
    for group in groups:
        p = content_frame.add_paragraph()
        p.text = group
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Add Ransomware Attack Lifecycle slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Ransomware Attack Lifecycle"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    lifecycle = [
        "1. Initial Access - Phishing, RDP, VPN exploits",
        "2. Reconnaissance - Network mapping and privilege escalation",
        "3. Lateral Movement - Spreading across the network",
        "4. Data Exfiltration - Stealing sensitive data",
        "5. Encryption - Encrypting files and systems",
        "6. Ransom Note - Demanding payment",
        "7. Negotiation - Communication with victims",
        "8. Payment/Extortion - Final ransom demands"
    ]
    
    for i, step in enumerate(lifecycle):
        p = content_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.space_before = Pt(6)
        if i == 0:
            p.font.bold = True
    
    # Add Ransomware Prevention slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Ransomware Prevention"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    prevention = [
        "• Regular offline backups (3-2-1 rule)",
        "• Employee security awareness training",
        "• Multi-factor authentication (MFA)",
        "• Patch management and vulnerability scanning",
        "• Network segmentation and zero trust",
        "• Endpoint detection and response (EDR)",
        "• Email filtering and anti-phishing",
        "• Incident response planning and testing"
    ]
    
    for item in prevention:
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Add VAJRA Ransomware Tracking slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "VAJRA Ransomware Tracking"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "VAJRA Platform provides comprehensive ransomware tracking capabilities:"
    
    p = content_frame.add_paragraph()
    p.text = "• Real-time ransomware incident tracking via Ransomware.live"
    p.font.size = Pt(18)
    p.space_before = Pt(12)
    p = content_frame.add_paragraph()
    p.text = "• Ransomware group monitoring and statistics"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• Geographic attack visualization"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• Industry-specific ransomware impact analysis"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• Ransomware trend analysis and forecasting"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• PDF report generation for ransomware incidents"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• Real-time alerts for new ransomware attacks"
    p.font.size = Pt(18)
    p = content_frame.add_paragraph()
    p.text = "• Historical ransomware data and incident analysis"
    p.font.size = Pt(18)
    
    # Save the updated presentation
    prs.save('DATE - 06-08-2026.pptx')
    print("Ransomware content added successfully to DATE - 06-08-2026.pptx")

if __name__ == "__main__":
    add_ransomware_content()
