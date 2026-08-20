import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_vajra_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(9, 13, 22)         # #090D16
    COLOR_CARD = RGBColor(17, 24, 39)      # #111827
    COLOR_BORDER = RGBColor(31, 41, 55)    # #1F2937
    COLOR_BLUE = RGBColor(59, 130, 246)    # #3B82F6
    COLOR_RED = RGBColor(239, 68, 68)      # #EF4444
    COLOR_AMBER = RGBColor(245, 158, 11)   # #F59E0B
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_GRAY = RGBColor(156, 163, 175)   # #9CA3AF

    def apply_dark_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG

    def add_header(slide, category_text, title_text):
        # Category Tag
        tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_BLUE
        
        # Main Title
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s1)
    
    # Hero Card Container
    add_card(s1, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), COLOR_CARD, COLOR_BLUE)
    
    tb1 = s1.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.733), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "🛡️ VAJRA THREAT INTELLIGENCE PLATFORM"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_BLUE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Ransomware Defense & Enterprise Risk Scoring"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf1.add_paragraph()
    p3.text = "\nReal-Time Telemetry Pipeline • Multi-Domain Perimeter Audit • Automated Executive PDF Briefings"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_GRAY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf1.add_paragraph()
    p4.text = "\nPrepared for CISO & Executive Cybersecurity Operations"
    p4.font.size = Pt(12)
    p4.font.color.rgb = COLOR_BLUE
    p4.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: Threat Landscape
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s2)
    add_header(s2, "EXECUTIVE CONTEXT", "The Escalating Global Ransomware Crisis")

    # 3 Stat Cards
    c1 = add_card(s2, Inches(0.8), Inches(1.6), Inches(3.6), Inches(5.2))
    tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "1,200+ DAILY ATTACKS"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_RED
    
    p = tf.add_paragraph()
    p.text = "\nHigh-Frequency Exploitation\n\nCybercriminal syndicates execute non-stop automated attacks targeting enterprise API endpoints and unpatched subdomains."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY

    c2 = add_card(s2, Inches(4.8), Inches(1.6), Inches(3.6), Inches(5.2))
    tb = s2.shapes.add_textbox(Inches(5.0), Inches(1.8), Inches(3.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "DOUBLE EXTORTION (RaaS)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    p = tf.add_paragraph()
    p.text = "\nData Exfiltration Prior to Encryption\n\nGroups like LockBit, BlackCat/ALPHV, and Cl0p steal sensitive enterprise IP before deploying locking ransomware."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY

    c3 = add_card(s2, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2))
    tb = s2.shapes.add_textbox(Inches(9.0), Inches(1.8), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "$1.85M AVERAGE DOWNTIME COST"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    p = tf.add_paragraph()
    p.text = "\nCritical Business Disruption\n\nExcludes reputational damage, legal compliance fines, and supply chain SLA breach penalties."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 3: Telemetry Architecture
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s3)
    add_header(s3, "PLATFORM ARCHITECTURE", "9 Parallel Telemetry Feeds Integrated into VAJRA Core")

    add_card(s3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = s3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "📡 INTEGRATED THREAT SOURCES"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    sources = [
        "1. NIST NVD API v2.0 (Live CVE & CVSS v3 Scores)",
        "2. Ransomware.live API (RaaS Victim Monitoring)",
        "3. AlienVault OTX (Open Threat Exchange Pulses)",
        "4. Live TLS Socket Inspection (Certifi CA Trust)",
        "5. ICANN RDAP WHOIS (Domain Age & Registrar)",
        "6. Google Public DNS REST API (DNS Records)",
        "7. AbuseIPDB API (Abuse Confidence Scoring)",
        "8. URLScan.io API (Malicious Script Scanning)",
        "9. VirusTotal API v3 (Multi-Antivirus Reputation)"
    ]
    for s in sources:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE

    add_card(s3, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = s3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "⚙️ VAJRA PARALLEL ANALYSIS PIPELINE"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_AMBER

    steps = [
        "• Real-Time Asynchronous Ingestion (Python Asyncio/HTTPX)",
        "• Unified 0-100 Security Risk Score Normalization",
        "• High/Critical Severity Alert Dispatching",
        "• 5-Tab Deep Domain Inspection Component Rendering",
        "• One-Click Automated PDF Executive Report Generation"
    ]
    for st in steps:
        p = tf.add_paragraph()
        p.text = f"\n{st}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 4: Ransomware Intelligence Focus
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s4)
    add_header(s4, "CORE FEATURE FOCUS 1", "Ransomware Intelligence & Real-Time Victim Tracking")

    add_card(s4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tb = s4.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    tf.paragraphs[0].text = "☣️ LIVE RANSOMWARE VICTIM RADAR"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_RED

    features = [
        ("Dark Web Leak Monitoring", "Monitors RaaS victim leak sites in real-time to alert security teams immediately when supply chain partners or targeted sectors appear."),
        ("Active Syndicate Profiling", "Tracks TTPs, target demographics, and attack frequency for groups like LockBit 3.0, BlackCat/ALPHV, Cl0p, Akira, and Play."),
        ("Early Warning Sector Radar", "Aggregates incident counts across Healthcare, Aviation, Finance, and Government to adjust defensive posture proactively."),
        ("Instant PDF Export", "Generates publication-ready Ransomware Intelligence Reports for CISO briefings and board compliance.")
    ]
    for title, desc in features:
        p = tf.add_paragraph()
        p.text = f"\n• {title}: "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.text += desc
        p.font.bold = False
        p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 5: Risk Score & Domain Inspection
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s5)
    add_header(s5, "CORE FEATURE FOCUS 2", "Automated 0-100 Risk Score & 5-Tab Inspection")

    add_card(s5, Inches(0.8), Inches(1.6), Inches(4.5), Inches(5.2))
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🧮 RISK SCORE ALGORITHM"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE

    p = tf.add_paragraph()
    p.text = "\n• Baseline Score: Starts at 100\n• Expired TLS Certificate: -30 points\n• AbuseIPDB Score > 50%: Score capped ≤50\n• URLScan Malicious Tag: Score capped ≤45\n• NVD Critical CVEs: Score penalized by CVSS v3 ratings"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_GRAY

    add_card(s5, Inches(5.6), Inches(1.6), Inches(6.9), Inches(5.2))
    tb = s5.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(6.5), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "🔍 5-TAB DEEP DOMAIN INSPECTION"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_AMBER

    tabs = [
        ("1. Overview & IPs", "Resolved IP addresses, Geolocation, ISP, Domain Age (days)"),
        ("2. Domain Threats", "Detected threats, severity breakdown, confidence ratings"),
        ("3. Vulnerabilities", "NVD CVE vulnerability list with direct NIST NVD links"),
        ("4. SSL / TLS Verification", "Live TLS handshake, CA trust chain, expiration countdown"),
        ("5. DNS & Feeds", "A / MX / TXT DNS records matrix")
    ]
    for tab_title, tab_desc in tabs:
        p = tf.add_paragraph()
        p.text = f"\n• {tab_title}: {tab_desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 6: Enterprise Use Case 1 (Aviation)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s6)
    add_header(s6, "ENTERPRISE USE CASE 1", "Aviation & Transportation Security (e.g., IndiGo Airlines)")

    add_card(s6, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    tf.paragraphs[0].text = "✈️ PROTECTING CRITICAL PASSENGER & BOOKING INFRASTRUCTURE"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE

    usecase1 = [
        ("Challenge", "Commercial airlines like IndiGo process millions of passenger transactions daily across dozens of subdomains and flight booking APIs, making them high-priority ransomware targets."),
        ("VAJRA Action", "Continuous parallel monitoring of enterprise domains ('goindigo.in'), auditing SSL certificate validity, open NVD CVEs, and AbuseIPDB indicator scores."),
        ("Key Result", "Detects expired subdomains, untrusted TLS certificates, and unpatched CVE vulnerabilities before cybercriminals launch ransomware or data exfiltration attacks.")
    ]
    for heading, text in usecase1:
        p = tf.add_paragraph()
        p.text = f"\n• {heading}: "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.text += text
        p.font.bold = False
        p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 7: Enterprise Use Case 2 (SOC Response)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s7)
    add_header(s7, "ENTERPRISE USE CASE 2", "SOC Incident Response & Automated Executive PDF Briefings")

    add_card(s7, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tb = s7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    tf.paragraphs[0].text = "⚡ ACCELERATING SOC INVESTIGATIONS & C-SUITE REPORTING"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_AMBER

    usecase2 = [
        ("Challenge", "Security operations analysts waste hours manually querying 8+ different intelligence APIs and formatting static reports for executive leadership during critical incidents."),
        ("VAJRA Action", "Analysts click 'View Details' to inspect domain security across separate browser tabs while leveraging VAJRA's one-click PDF Report engine."),
        ("Key Result", "90% reduction in investigation time. Instant generation of publication-ready PDF Executive Summaries, Ransomware Reports, and Company Risk Assessments directly from the top navigation bar.")
    ]
    for heading, text in usecase2:
        p = tf.add_paragraph()
        p.text = f"\n• {heading}: "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.text += text
        p.font.bold = False
        p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 8: CISO Key Benefits
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s8)
    add_header(s8, "STRATEGIC ROI", "Key Benefits for CISO & Executive Leadership")

    # 4 Benefit Cards Grid
    b_cards = [
        ("🛡️ PROACTIVE RANSOMWARE DEFENSE", "Stops initial access by fixing expired TLS certificates, high AbuseIPDB scores, and critical NVD CVEs before encryption happens.", COLOR_BLUE),
        ("⏱️ 90% ANALYST TIME SAVINGS", "Consolidates 9 separate threat intelligence feeds into a single automated dashboard, eliminating manual tool hopping.", COLOR_AMBER),
        ("📊 C-SUITE BOARDROOM VISIBILITY", "Standardized 0-100 Risk Scores and instant one-click PDF exports simplify security presentation to executive boards.", COLOR_BLUE),
        ("🔗 SUPPLY CHAIN PROTECTION", "Continuously audits third-party vendor subdomains and partner infrastructure to prevent supply chain compromise.", COLOR_RED)
    ]

    coords = [
        (Inches(0.8), Inches(1.6)), (Inches(6.8), Inches(1.6)),
        (Inches(0.8), Inches(4.3)), (Inches(6.8), Inches(4.3))
    ]

    for idx, (b_title, b_desc, b_color) in enumerate(b_cards):
        left, top = coords[idx]
        add_card(s8, left, top, Inches(5.7), Inches(2.4))
        tb = s8.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = b_title
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = b_color

        p = tf.add_paragraph()
        p.text = f"\n{b_desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_GRAY

    # ==========================================
    # SLIDE 9: Feature Matrix
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s9)
    add_header(s9, "CAPABILITIES MATRIX", "VAJRA Platform Core Feature Summary")

    add_card(s9, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tb = s9.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True

    matrix = [
        ("Ransomware Live Radar", "Monitors active RaaS leak sites & victim announcements", "Early warning threat radar"),
        ("NVD CVE Integration", "Maps CVE vulnerabilities & CVSS v3 base scores", "Rapid patch prioritization"),
        ("5-Tab Domain Inspection", "IP, Threats, CVEs, SSL Audit, DNS Matrix", "360-degree perimeter visibility"),
        ("Live TLS Socket Verification", "Direct CA trust store handshake via Certifi", "Prevents certificate outages"),
        ("One-Click PDF Export", "Executive, Technical, & Combined PDF Reports", "Effortless compliance reporting")
    ]

    tf.paragraphs[0].text = "FEATURE                                        DESCRIPTION                                                             STRATEGIC VALUE"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE

    for f_name, f_desc, f_val in matrix:
        p = tf.add_paragraph()
        p.text = f"\n• {f_name.ljust(25)} : {f_desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 10: Conclusion & Roadmap
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    apply_dark_bg(s10)
    add_header(s10, "STRATEGIC ROADMAP", "Conclusion & Next Steps")

    add_card(s10, Inches(1.5), Inches(1.6), Inches(10.333), Inches(5.2), COLOR_CARD, COLOR_BLUE)
    tb = s10.shapes.add_textbox(Inches(1.8), Inches(1.9), Inches(9.733), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True

    tf.paragraphs[0].text = "🚀 ENTERPRISE DEPLOYMENT ROADMAP"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_BLUE

    roadmap = [
        ("Phase 1: Perimeter Asset Onboarding", "Register core enterprise domains, subdomains, and key third-party vendor assets in VAJRA Company Monitoring."),
        ("Phase 2: SOC Threshold Configuration", "Configure automated high-severity alert thresholds and real-time WebSocket notifications."),
        ("Phase 3: Executive Reporting Integration", "Integrate automated one-click PDF reports into monthly C-Suite security reviews and board presentations.")
    ]
    for phase_title, phase_desc in roadmap:
        p = tf.add_paragraph()
        p.text = f"\n• {phase_title}: "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.text += phase_desc
        p.font.bold = False
        p.font.color.rgb = COLOR_GRAY

    p_end = tf.add_paragraph()
    p_end.text = "\nThank you. Open Floor for Executive Q&A."
    p_end.font.size = Pt(14)
    p_end.font.bold = True
    p_end.font.color.rgb = COLOR_AMBER
    p_end.alignment = PP_ALIGN.CENTER

    # Save Presentation to File System
    artifact_path = "/Users/surajmujumdar/.gemini/antigravity-ide/brain/78a87757-f334-4d6c-a21d-c1e6a6ecfbec/VAJRA_Threat_Intelligence_Ransomware.pptx"
    workspace_path = "/Users/surajmujumdar/Desktop/indigo/new/INDIGO/VAJRA_Threat_Intelligence_Ransomware.pptx"
    public_path = "/Users/surajmujumdar/Desktop/indigo/new/INDIGO/frontend/public/VAJRA_Threat_Intelligence_Ransomware.pptx"

    prs.save(artifact_path)
    prs.save(workspace_path)
    prs.save(public_path)
    print("✅ PowerPoint presentation created successfully!")

if __name__ == "__main__":
    create_vajra_ppt()
