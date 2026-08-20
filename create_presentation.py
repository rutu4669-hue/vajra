#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "VAJRA Security Platform"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Complete Workflow & Architecture Overview"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Platform Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Platform Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "VAJRA Security Platform is a comprehensive cybersecurity threat intelligence and monitoring system designed for Security Operations Centers (SOCs)."
    
    p = content_frame.add_paragraph()
    p.text = "Key Features:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    
    features = [
        "• Real-time threat intelligence from AlienVault OTX",
        "• Ransomware attack tracking via Ransomware.live",
        "• Company risk assessment and monitoring",
        "• Live cyber threat news from Hacker News",
        "• PDF report generation with black-and-white styling",
        "• AI-powered security analysis via Gemini",
        "• WebSocket real-time updates",
        "• Comprehensive audit logging"
    ]
    
    for feature in features:
        p = content_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.level = 1
        p.space_before = Pt(6)
    
    # Slide 3: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "Three-Tier Architecture:"
    
    p = content_frame.add_paragraph()
    p.text = "Frontend Layer (Next.js)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p = content_frame.add_paragraph()
    p.text = "• React-based dashboard with real-time updates"
    p.font.size = Pt(16)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "Backend Layer (FastAPI)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p = content_frame.add_paragraph()
    p.text = "• RESTful API with 13 modules and 30+ endpoints"
    p.font.size = Pt(16)
    p.level = 1
    p = content_frame.add_paragraph()
    p.text = "• JWT authentication and authorization"
    p.font.size = Pt(16)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "Data Layer (PostgreSQL)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    p = content_frame.add_paragraph()
    p.text = "• 8 database models with relationships"
    p.font.size = Pt(16)
    p.level = 1
    p = content_frame.add_paragraph()
    p.text = "• Comprehensive audit logging"
    p.font.size = Pt(16)
    p.level = 1
    
    # Slide 4: User Journey Flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "User Journey Flow"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    steps = [
        "1. User Login → JWT Token Generation",
        "2. Dashboard Load → Real-time Data Fetch",
        "3. Module Navigation → Threat Intelligence, Ransomware, News, etc.",
        "4. Company Monitoring → Risk Assessment & Threat Analysis",
        "5. Report Generation → PDF Download",
        "6. Real-time Updates → WebSocket Notifications",
        "7. Activity Logging → Audit Trail",
        "8. Logout → Token Clearance"
    ]
    
    for i, step in enumerate(steps):
        p = content_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.space_before = Pt(8)
        if i == 0:
            p.font.bold = True
    
    # Slide 5: Database Models
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Database Models"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    models = [
        "User Model - Authentication & user management",
        "Company Model - Organization monitoring",
        "CompanyThreat Model - Threat tracking per company",
        "CompanyRiskAssessment Model - Risk scoring",
        "Alert Model - Security alerts",
        "RansomwareIncident Model - Ransomware attack tracking",
        "AttackEvent Model - Geographic attack mapping",
        "ActivityLog Model - Complete audit trail"
    ]
    
    for model in models:
        p = content_frame.add_paragraph()
        p.text = model
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Slide 6: External API Integrations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "External API Integrations"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    apis = [
        "AlienVault OTX - Threat intelligence & IOC data",
        "Ransomware.live - Ransomware attack tracking",
        "Hacker News - Cyber threat news feed",
        "VirusTotal - Malware analysis",
        "AbuseIPDB - IP reputation checking",
        "SSL Labs - SSL certificate analysis",
        "Domain Analysis APIs - Domain risk assessment",
        "Gemini AI - AI-powered security analysis",
        "HuggingFace - Machine learning models"
    ]
    
    for api in apis:
        p = content_frame.add_paragraph()
        p.text = api
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Slide 7: API Modules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "API Modules"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    modules = [
        "Authentication - User login, registration, token management",
        "Dashboard - Summary, alerts, attack map data",
        "Threat Intelligence - Metrics, trends, actors, industries",
        "Ransomware - Incidents, statistics, group tracking",
        "News - Live cyber threat news",
        "Reports - PDF generation for all modules",
        "AI - AI-powered analysis and predictions",
        "Companies - Risk assessment and monitoring",
        "Domain Analysis - Domain risk evaluation",
        "Alerts - Security alert management",
        "WebSocket - Real-time data streaming"
    ]
    
    for module in modules:
        p = content_frame.add_paragraph()
        p.text = module
        p.font.size = Pt(15)
        p.space_before = Pt(5)
    
    # Slide 8: Security Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Security Features"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    security = [
        "JWT Authentication - Token-based access control",
        "Password Hashing - Secure password storage",
        "CORS - Cross-origin resource sharing",
        "Rate Limiting - API abuse prevention",
        "Input Validation - SQL injection & XSS prevention",
        "Security Headers - HTTP security best practices",
        "Activity Logging - Complete audit trail",
        "Role-based Access Control - User permission management"
    ]
    
    for feature in security:
        p = content_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Slide 9: Report Generation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Report Generation"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    reports = [
        "Overall PDF Report - Comprehensive platform overview",
        "Executive Summary - High-level security metrics",
        "Threat Intelligence Report - Detailed threat analysis",
        "Ransomware Report - Ransomware incident tracking",
        "Company Report - Organization-specific risk assessment",
        "Black-and-White Styling - Consistent professional appearance",
        "Real-time Generation - On-demand PDF creation",
        "Download Options - Multiple report formats"
    ]
    
    for report in reports:
        p = content_frame.add_paragraph()
        p.text = report
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Slide 10: Real-time Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Real-time Features"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    
    realtime = [
        "WebSocket Connection - Persistent bidirectional communication",
        "Live Threat Updates - Real-time threat intelligence",
        "Alert Notifications - Instant security alerts",
        "Dashboard Refresh - Automatic data updates",
        "Background Sync - Scheduled data synchronization",
        "Cache Management - Redis caching for performance",
        "Event Broadcasting - Multi-client updates",
        "Low Latency - Sub-second response times"
    ]
    
    for feature in realtime:
        p = content_frame.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Slide 11: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "VAJRA Security Platform provides:"
    
    p = content_frame.add_paragraph()
    p.text = "• Comprehensive threat intelligence and monitoring"
    p.font.size = Pt(20)
    p.space_before = Pt(12)
    p = content_frame.add_paragraph()
    p.text = "• Real-time data updates and notifications"
    p.font.size = Pt(20)
    p = content_frame.add_paragraph()
    p.text = "• Professional PDF reporting with consistent styling"
    p.font.size = Pt(20)
    p = content_frame.add_paragraph()
    p.text = "• AI-powered security analysis"
    p.font.size = Pt(20)
    p = content_frame.add_paragraph()
    p.text = "• Complete audit trail and activity logging"
    p.font.size = Pt(20)
    p = content_frame.add_paragraph()
    p.text = "• Scalable architecture for enterprise deployment"
    p.font.size = Pt(20)
    
    p = content_frame.add_paragraph()
    p.text = "Built for Security Operations Centers to enhance threat detection and response capabilities."
    p.font.size = Pt(18)
    p.font.italic = True
    p.space_before = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    prs.save('VAJRA_Platform_Workflow.pptx')
    print("PowerPoint presentation created successfully: VAJRA_Platform_Workflow.pptx")

if __name__ == "__main__":
    create_presentation()
